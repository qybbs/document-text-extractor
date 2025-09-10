from fastapi import FastAPI, File, UploadFile, HTTPException
import pytesseract
from PIL import Image
import pdfplumber
import io
import os
from docx import Document
import openpyxl
from pptx import Presentation

app = FastAPI(title="Document Text Extractor API")

pytesseract.pytesseract.tesseract_cmd = os.getenv("TESSERACT_CMD", "tesseract")


@app.get("/health")
async def health():
    """Simple health check endpoint"""
    return {"status": "ok"}


@app.post("/extract/pdf")
async def extract_pdf(file: UploadFile = File(...)):
    if not file.filename.lower().endswith(".pdf"):
        raise HTTPException(status_code=400, detail="Invalid file format, expected PDF")

    text_result = []
    try:
        pdf_bytes = await file.read()
        with pdfplumber.open(io.BytesIO(pdf_bytes)) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_result.append(page_text)
                else:
                    im = page.to_image(resolution=300).original
                    ocr_text = pytesseract.image_to_string(im, lang="eng")
                    text_result.append(ocr_text)
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing PDF: {str(e)}")

    result = "\n".join(text_result).strip()
    return {"text": result}


@app.post("/extract/image")
async def extract_image(file: UploadFile = File(...)):
    if file.content_type not in [
        "image/jpeg", "image/png", "image/jpg", "image/gif", "image/webp"
    ]:
        raise HTTPException(status_code=400, detail="Invalid image format")

    try:
        img_bytes = await file.read()
        img = Image.open(io.BytesIO(img_bytes))
        text = pytesseract.image_to_string(img, lang="eng")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing image: {str(e)}")

    return {"text": text.strip()}


@app.post("/extract/office")
async def extract_office(file: UploadFile = File(...)):
    ext = file.filename.split(".")[-1].lower()
    file_bytes = await file.read()
    buffer = io.BytesIO(file_bytes)

    try:
        if ext == "docx":
            text = extract_docx(buffer)
        elif ext == "xlsx":
            text = extract_xlsx(buffer)
        elif ext == "pptx":
            text = extract_pptx(buffer)
        else:
            raise HTTPException(status_code=400, detail="Unsupported office format")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error processing office file: {str(e)}")

    return {"text": text.strip()}


def extract_docx(file_obj: io.BytesIO) -> str:
    """Extract text from DOCX"""
    doc = Document(file_obj)
    text = [para.text for para in doc.paragraphs if para.text.strip()]
    return "\n".join(text)


def extract_xlsx(file_obj: io.BytesIO) -> str:
    """Extract text from XLSX (all cells)"""
    wb = openpyxl.load_workbook(file_obj, data_only=True)
    text = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        for row in ws.iter_rows(values_only=True):
            for cell in row:
                if cell:
                    text.append(str(cell))
    return "\n".join(text)


def extract_pptx(file_obj: io.BytesIO) -> str:
    """Extract text from PPTX (all slides)"""
    prs = Presentation(file_obj)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text.append(shape.text)
    return "\n".join(text)
